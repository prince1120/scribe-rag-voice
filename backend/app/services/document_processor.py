import os
from typing import List, Dict, Any, Optional
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd
import numpy as np
from pathlib import Path
import logging
from uuid import uuid4
import re

try:
    from PIL import Image
    _PIL_AVAILABLE = True
except ImportError:
    _PIL_AVAILABLE = False

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

_SENTENCE_SPLIT_RE = re.compile(
    r'(?<=[.!?])\s+(?=[A-Z0-9"‘“(])|\n\s*\n'
)


class DocumentProcessor:
    """Handle multiple document formats."""

    def __init__(
        self,
        chunk_size: int = 512,
        chunk_overlap: int = 50,
        vision_ocr=None,
        embedding_service=None,
        semantic_chunking: bool = True,
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        # vision_ocr: optional VisionOCR instance (Groq vision model).
        # Used for: image uploads + scanned-PDF page OCR fallback.
        self.vision_ocr = vision_ocr
        # embedding_service: optional EmbeddingService (bi-encoder). When
        # present, text is split into semantically coherent chunks instead
        # of naive fixed-word windows — see `_semantic_chunk_text`.
        self.embedding_service = embedding_service if semantic_chunking else None
    
    def process_file(self, file_path: str, metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Process a file and return chunks with metadata."""
        file_ext = Path(file_path).suffix.lower()
        
        processors = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.pptx': self._process_pptx,
            '.txt': self._process_text,
            '.md': self._process_text,
            '.html': self._process_text,
            '.csv': self._process_csv,
            '.xlsx': self._process_excel,
            '.png': self._process_image,
            '.jpg': self._process_image,
            '.jpeg': self._process_image,
            '.webp': self._process_image,
            '.bmp': self._process_image,
            '.tiff': self._process_image,
            '.tif': self._process_image,
            '.gif': self._process_image,
        }
        
        if file_ext not in processors:
            raise ValueError(f"Unsupported file type: {file_ext}")
        
        return processors[file_ext](file_path, metadata)
    
    def _process_pdf(self, file_path: str, metadata: Dict) -> List[Dict]:
        """Process PDF files. Falls back to per-page OCR for scanned/image PDFs."""
        chunks = []
        try:
            reader = PdfReader(file_path)
            full_text = []

            for page_num, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    full_text.append((page_num + 1, text))

            total_extracted = sum(len(t) for _, t in full_text)
            page_count = len(reader.pages)

            # Likely a scanned / slide-image PDF — fall back to OCR.
            if total_extracted < max(50, page_count * 20):
                logger.info(
                    f"PDF text extraction yielded only {total_extracted} chars across "
                    f"{page_count} pages — running OCR fallback."
                )
                full_text = self._ocr_pdf_pages(file_path)

            for page_num, text in full_text:
                page_chunks = self._chunk(text)
                for chunk_text in page_chunks:
                    chunks.append({
                        "chunk_id": str(uuid4()),
                        "document_id": metadata["document_id"],
                        "content": chunk_text,
                        "chunk_index": len(chunks),
                        "page_number": page_num,
                        "metadata": {**metadata, "page_number": page_num},
                    })

            logger.info(f"Processed PDF: {len(chunks)} chunks ({page_count} pages)")
        except Exception as e:
            logger.error(f"Error processing PDF: {e}")
            raise

        return chunks

    def _ocr_pdf_pages(self, file_path: str) -> List[tuple]:
        """Rasterize each PDF page and OCR it via Groq vision model.

        Returns [(page_number, text), ...]. Raises if no OCR backend is
        configured.
        """
        if self.vision_ocr is None:
            raise RuntimeError(
                "Scanned PDF needs an OCR backend but none is configured. "
                "Pass a VisionOCR instance to DocumentProcessor."
            )
        try:
            import pypdfium2 as pdfium
        except ImportError as e:
            raise RuntimeError(
                "Scanned PDF support requires `pypdfium2`. "
                "Install with: pip install pypdfium2"
            ) from e

        pages_text: List[tuple] = []
        pdf = pdfium.PdfDocument(file_path)
        try:
            for idx in range(len(pdf)):
                page = pdf[idx]
                # scale=2 ≈ 144 DPI — good legibility for vision model
                bitmap = page.render(scale=2.0)
                pil_image = bitmap.to_pil()
                try:
                    text = self.vision_ocr.extract_text(pil_image) or ""
                except Exception as e:
                    logger.warning(f"Vision OCR failed on page {idx + 1}: {e}")
                    text = ""
                text = re.sub(r"\s+\n", "\n", text).strip()
                if text:
                    pages_text.append((idx + 1, text))
                logger.info(f"OCR page {idx + 1}/{len(pdf)}: {len(text)} chars")
        finally:
            pdf.close()
        return pages_text
    
    def _process_docx(self, file_path: str, metadata: Dict) -> List[Dict]:
        """Process Word documents."""
        chunks = []
        try:
            doc = DocxDocument(file_path)
            full_text = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            # Also extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    full_text.append(row_text)
            
            combined_text = "\n".join(full_text)
            text_chunks = self._chunk(combined_text)
            
            for i, chunk_text in enumerate(text_chunks):
                chunks.append({
                    "chunk_id": str(uuid4()),
                    "document_id": metadata["document_id"],
                    "content": chunk_text,
                    "chunk_index": i,
                    "metadata": metadata
                })
            
            logger.info(f"Processed DOCX: {len(chunks)} chunks")
        except Exception as e:
            logger.error(f"Error processing DOCX: {e}")
            raise
        
        return chunks
    
    def _process_pptx(self, file_path: str, metadata: Dict) -> List[Dict]:
        """Process PowerPoint files."""
        chunks = []
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    combined = "\n".join(slide_text)
                    text_chunks = self._chunk(combined)
                    
                    for i, chunk_text in enumerate(text_chunks):
                        chunks.append({
                            "chunk_id": str(uuid4()),
                            "document_id": metadata["document_id"],
                            "content": chunk_text,
                            "chunk_index": len(chunks),
                            "slide_number": slide_num + 1,
                            "metadata": {**metadata, "slide_number": slide_num + 1}
                        })
            
            logger.info(f"Processed PPTX: {len(chunks)} chunks")
        except Exception as e:
            logger.error(f"Error processing PPTX: {e}")
            raise
        
        return chunks
    
    def _process_text(self, file_path: str, metadata: Dict) -> List[Dict]:
        """Process plain text, Markdown, or HTML files."""
        chunks = []
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()

            text_chunks = self._chunk(text)
            
            for i, chunk_text in enumerate(text_chunks):
                chunks.append({
                    "chunk_id": str(uuid4()),
                    "document_id": metadata["document_id"],
                    "content": chunk_text,
                    "chunk_index": i,
                    "metadata": metadata
                })
            
            logger.info(f"Processed text file: {len(chunks)} chunks")
        except Exception as e:
            logger.error(f"Error processing text file: {e}")
            raise
        
        return chunks
    
    # Safety: a 50 MB CSV could hold hundreds of thousands of rows — one chunk
    # per row would OOM during embedding and upsert. Cap rows and group them
    # so the output stays bounded regardless of file size.
    _MAX_CSV_ROWS = 5000
    _CSV_ROWS_PER_CHUNK = 20

    def _process_csv(self, file_path: str, metadata: Dict) -> List[Dict]:
        """Process CSV files — bounded: groups rows and caps total rows."""
        chunks: List[Dict] = []
        try:
            df = pd.read_csv(file_path)

            if len(df) > self._MAX_CSV_ROWS:
                logger.warning(
                    f"CSV has {len(df)} rows; capping at {self._MAX_CSV_ROWS} "
                    f"({self._CSV_ROWS_PER_CHUNK} rows per chunk)"
                )
                df = df.head(self._MAX_CSV_ROWS)

            for start in range(0, len(df), self._CSV_ROWS_PER_CHUNK):
                batch = df.iloc[start : start + self._CSV_ROWS_PER_CHUNK]
                lines = []
                for _, row in batch.iterrows():
                    lines.append(" | ".join(f"{col}: {val}" for col, val in row.items()))
                row_text = "\n".join(lines)
                chunk_index = len(chunks)
                chunks.append({
                    "chunk_id": str(uuid4()),
                    "document_id": metadata["document_id"],
                    "content": row_text,
                    "chunk_index": chunk_index,
                    "row_number": int(batch.index[0]) + 1,
                    "metadata": {**metadata, "row_number": int(batch.index[0]) + 1},
                })

            logger.info(f"Processed CSV: {len(chunks)} chunks from {len(df)} rows")
        except Exception as e:
            logger.error(f"Error processing CSV: {e}")
            raise

        return chunks
    
    _MAX_EXCEL_ROWS_TOTAL = 5000
    _EXCEL_ROWS_PER_CHUNK = 20

    def _process_excel(self, file_path: str, metadata: Dict) -> List[Dict]:
        """Process Excel files — bounded across all sheets."""
        chunks: List[Dict] = []
        try:
            xl_file = pd.ExcelFile(file_path)

            rows_seen = 0
            for sheet_name in xl_file.sheet_names:
                if rows_seen >= self._MAX_EXCEL_ROWS_TOTAL:
                    logger.warning(
                        f"Excel capped at {self._MAX_EXCEL_ROWS_TOTAL} rows total; "
                        f"skipping remaining sheets after '{sheet_name}'"
                    )
                    break
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                remaining = self._MAX_EXCEL_ROWS_TOTAL - rows_seen
                if len(df) > remaining:
                    logger.warning(f"Sheet '{sheet_name}' capped from {len(df)} to {remaining} rows")
                    df = df.head(remaining)

                sheet_metadata = {**metadata, "sheet_name": sheet_name}
                for start in range(0, len(df), self._EXCEL_ROWS_PER_CHUNK):
                    batch = df.iloc[start : start + self._EXCEL_ROWS_PER_CHUNK]
                    lines = []
                    for _, row in batch.iterrows():
                        lines.append(" | ".join(f"{col}: {val}" for col, val in row.items()))
                    row_text = f"Sheet: {sheet_name} | " + "\n".join(lines)
                    chunks.append({
                        "chunk_id": str(uuid4()),
                        "document_id": metadata["document_id"],
                        "content": row_text,
                        "chunk_index": len(chunks),
                        "row_number": int(batch.index[0]) + 1,
                        "sheet_name": sheet_name,
                        "metadata": sheet_metadata,
                    })
                rows_seen += len(df)

            logger.info(f"Processed Excel: {len(chunks)} chunks from {rows_seen} rows")
        except Exception as e:
            logger.error(f"Error processing Excel: {e}")
            raise
        
        return chunks
    
    def _process_image(self, file_path: str, metadata: Dict) -> List[Dict]:
        """Process image files via Groq vision-model OCR.

        Extracts any text visible in the image so the image becomes searchable
        alongside other documents. Always emits at least one chunk so the
        document remains discoverable even if OCR finds nothing useful.
        """
        if not _PIL_AVAILABLE:
            raise RuntimeError("Pillow is required to process images.")
        if self.vision_ocr is None:
            raise RuntimeError(
                "Image upload needs an OCR backend but none is configured. "
                "Pass a VisionOCR instance to DocumentProcessor."
            )

        chunks: List[Dict] = []
        filename = metadata.get("filename", os.path.basename(file_path))
        try:
            with Image.open(file_path) as img:
                width, height = img.size
            try:
                raw_text = self.vision_ocr.extract_text(file_path) or ""
            except Exception as e:
                logger.warning(f"Vision OCR failed for image: {e}")
                raw_text = ""

            ocr_text = re.sub(r"\s+\n", "\n", raw_text).strip()

            image_meta = {
                **metadata,
                "is_image": True,
                "file_path": os.path.abspath(file_path),
                "image_width": width,
                "image_height": height,
                "ocr_char_count": len(ocr_text),
            }

            if ocr_text:
                text_chunks = self._chunk(ocr_text)
                for i, chunk_text in enumerate(text_chunks):
                    body = f"[Image: {filename}]\n{chunk_text}"
                    chunks.append({
                        "chunk_id": str(uuid4()),
                        "document_id": metadata["document_id"],
                        "content": body,
                        "chunk_index": i,
                        "metadata": image_meta,
                    })
            else:
                # No OCR text — still emit a stub so retrieval can find the image by filename
                chunks.append({
                    "chunk_id": str(uuid4()),
                    "document_id": metadata["document_id"],
                    "content": f"[Image: {filename}] (no readable text)",
                    "chunk_index": 0,
                    "metadata": image_meta,
                })

            logger.info(
                f"Processed image '{filename}': {len(chunks)} chunks, "
                f"{len(ocr_text)} OCR chars"
            )
        except Exception as e:
            logger.error(f"Error processing image: {e}")
            raise

        return chunks

    def _chunk_text(self, text: str) -> List[str]:
        """Split text into overlapping fixed-size chunks (word count)."""
        words = text.split()
        chunks = []

        for i in range(0, len(words), self.chunk_size - self.chunk_overlap):
            chunk_words = words[i:i + self.chunk_size]
            chunks.append(" ".join(chunk_words))

            if i + self.chunk_size >= len(words):
                break

        return chunks if chunks else [text]

    def chunk_text(self, text: str) -> List[str]:
        """Public wrapper for `_chunk` — used when content is already plain
        text and there's no file to run through `process_file` (e.g.
        re-indexing an edited image's OCR text)."""
        return self._chunk(text)

    def _chunk(self, text: str) -> List[str]:
        """Chunk text, preferring semantic (meaning-based) boundaries.

        Falls back to fixed-size word windows if no embedding service is
        configured or semantic chunking fails for any reason.
        """
        if self.embedding_service is not None:
            try:
                return self._semantic_chunk_text(text)
            except Exception as e:
                logger.warning(f"Semantic chunking failed, falling back to fixed-size: {e}")
        return self._chunk_text(text)

    def _split_sentences(self, text: str) -> List[str]:
        """Split text into sentences, treating paragraph breaks as hard boundaries."""
        text = text.strip()
        if not text:
            return []
        parts = _SENTENCE_SPLIT_RE.split(text)
        return [p.strip() for p in parts if p and p.strip()]

    def _semantic_chunk_text(
        self,
        text: str,
        percentile_threshold: float = 90.0,
        buffer_size: int = 1,
    ) -> List[str]:
        """Split text at semantic breakpoints (Kamradt-style percentile method).

        Each sentence is embedded together with its immediate neighbors for a
        more stable representation, then cosine *distance* between
        consecutive sentence embeddings is computed. Distances above the
        given percentile mark a topic shift, and the text is split there —
        so a chunk stays a single coherent idea instead of a mid-sentence,
        mid-thought cut. Oversized resulting chunks are still hard-split by
        word count so we never blow the retrieval/context token budget, and
        undersized ones are merged forward to avoid noisy tiny chunks.
        """
        sentences = self._split_sentences(text)
        if len(sentences) < 4:
            # Not enough structure to find meaningful semantic breakpoints.
            return self._chunk_text(text)

        combined = []
        for i in range(len(sentences)):
            start = max(0, i - buffer_size)
            end = min(len(sentences), i + buffer_size + 1)
            combined.append(" ".join(sentences[start:end]))

        embeddings = np.asarray(self.embedding_service.encode(combined))
        norms = np.linalg.norm(embeddings, axis=1, keepdims=True)
        normed = embeddings / np.clip(norms, 1e-8, None)

        # Cosine distance between consecutive (neighbor-buffered) sentences
        sims = np.sum(normed[:-1] * normed[1:], axis=1)
        distances = 1.0 - sims
        if distances.size == 0:
            return self._chunk_text(text)

        threshold = float(np.percentile(distances, percentile_threshold))
        breakpoints = {i for i, d in enumerate(distances) if d > threshold}

        raw_chunks: List[str] = []
        current = [sentences[0]]
        for i in range(1, len(sentences)):
            if (i - 1) in breakpoints:
                raw_chunks.append(" ".join(current))
                current = [sentences[i]]
            else:
                current.append(sentences[i])
        if current:
            raw_chunks.append(" ".join(current))

        # Merge undersized chunks forward, hard-split oversized ones.
        min_words = max(20, self.chunk_size // 6)
        final_chunks: List[str] = []
        pending = ""
        for chunk in raw_chunks:
            candidate = f"{pending} {chunk}".strip() if pending else chunk
            if len(candidate.split()) < min_words:
                pending = candidate
                continue
            pending = ""
            if len(candidate.split()) > self.chunk_size:
                final_chunks.extend(self._chunk_text(candidate))
            else:
                final_chunks.append(candidate)
        if pending:
            if len(pending.split()) > self.chunk_size:
                final_chunks.extend(self._chunk_text(pending))
            else:
                final_chunks.append(pending)

        return final_chunks if final_chunks else self._chunk_text(text)
