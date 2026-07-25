"""Extract editable plain text from a document, and regenerate the file
from edited text on save.

Scope (agreed tradeoff): this is *lossy* for binary office formats — a
round-tripped .docx/.pptx/.xlsx/.pdf keeps the text content but not the
original fonts, images, or layout. True fidelity editing would need an
external engine (OnlyOffice/Collabora); this is the pragmatic alternative
for in-app text editing without that infrastructure.
"""
import io
import re
from typing import List

import pandas as pd
from docx import Document as DocxDocument
from fpdf import FPDF
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches

EDITABLE_EXTENSIONS = {
    ".txt", ".md", ".html", ".csv", ".docx", ".pptx", ".xlsx", ".pdf",
}

_SLIDE_HEADER_RE = re.compile(r"^## Slide \d+\s*$", re.MULTILINE)
_PAGE_HEADER_RE = re.compile(r"^## Page \d+\s*$", re.MULTILINE)
_SHEET_HEADER_RE = re.compile(r"^## Sheet: (.+)$", re.MULTILINE)


def is_editable(ext: str) -> bool:
    return ext.lower() in EDITABLE_EXTENSIONS


def extract_editable_text(file_path: str, ext: str) -> str:
    """Return a plain-text representation of the file's content, suitable
    for display in a textarea and for re-parsing by write_editable_text."""
    ext = ext.lower()

    if ext in (".txt", ".md", ".html", ".csv"):
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    if ext == ".docx":
        doc = DocxDocument(file_path)
        return "\n\n".join(p.text for p in doc.paragraphs)

    if ext == ".pptx":
        prs = Presentation(file_path)
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = [
                shape.text for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            parts.append(f"## Slide {i}\n" + "\n".join(texts))
        return "\n\n".join(parts)

    if ext == ".xlsx":
        xl = pd.ExcelFile(file_path)
        parts = []
        for sheet in xl.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet)
            parts.append(f"## Sheet: {sheet}\n{df.to_csv(index=False)}")
        return "\n\n".join(parts)

    if ext == ".pdf":
        reader = PdfReader(file_path)
        parts = []
        for i, page in enumerate(reader.pages, start=1):
            parts.append(f"## Page {i}\n{page.extract_text() or ''}")
        return "\n\n".join(parts)

    raise ValueError(f"Editing not supported for '{ext}'")


def write_editable_text(file_path: str, ext: str, text: str) -> None:
    """Regenerate the file at file_path from edited plain text."""
    ext = ext.lower()

    if ext in (".txt", ".md", ".html", ".csv"):
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(text)
        return

    if ext == ".docx":
        doc = DocxDocument()
        for para in text.split("\n\n"):
            doc.add_paragraph(para)
        doc.save(file_path)
        return

    if ext == ".pptx":
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        sections = [s.strip() for s in _SLIDE_HEADER_RE.split(text) if s.strip()]
        if not sections:
            sections = [text.strip()] if text.strip() else [""]
        for section in sections:
            slide = prs.slides.add_slide(blank_layout)
            box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
            box.text_frame.text = section
        prs.save(file_path)
        return

    if ext == ".xlsx":
        sheets: List[tuple] = []
        headers = list(_SHEET_HEADER_RE.finditer(text))
        if headers:
            for idx, m in enumerate(headers):
                name = m.group(1).strip()[:31] or f"Sheet{idx + 1}"
                start = m.end()
                end = headers[idx + 1].start() if idx + 1 < len(headers) else len(text)
                csv_text = text[start:end].strip()
                sheets.append((name, csv_text))
        else:
            sheets = [("Sheet1", text.strip())]

        with pd.ExcelWriter(file_path, engine="openpyxl") as writer:
            for name, csv_text in sheets:
                try:
                    df = pd.read_csv(io.StringIO(csv_text))
                except Exception:
                    df = pd.DataFrame({"content": csv_text.splitlines()})
                df.to_excel(writer, sheet_name=name, index=False)
        return

    if ext == ".pdf":
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        sections = [s.strip() for s in _PAGE_HEADER_RE.split(text) if s.strip()]
        if not sections:
            sections = [text.strip()] if text.strip() else [""]
        for section in sections:
            pdf.add_page()
            pdf.set_font("Helvetica", size=11)
            # Core PDF fonts are latin-1 only; replace anything outside that
            # range rather than crashing on save.
            safe_text = section.encode("latin-1", errors="replace").decode("latin-1")
            pdf.multi_cell(0, 6, safe_text)
        pdf.output(file_path)
        return

    raise ValueError(f"Editing not supported for '{ext}'")
